from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.shapes.graphfrm import GraphicFrame

import io
import json

from pptx.shapes.group import GroupShape
from pptx.shapes.graphfrm import GraphicFrame
from pptx.chart.chart import Chart
from pptx.enum.chart import XL_CHART_TYPE
        
def find_chart(shape):
    if isinstance(shape, GraphicFrame) and shape.has_chart:
        return shape.chart
    elif isinstance(shape, GroupShape):
        for sub_shape in shape.shapes:
            chart = find_chart(sub_shape)
            if chart:
                return chart
    return None

def find_charts_in_slide(slide):
    charts = []

    def traverse(shape):
        if isinstance(shape, GraphicFrame) and shape.has_chart:
            charts.append(shape.chart)
        elif isinstance(shape, GroupShape):
            for sub_shape in shape.shapes:
                traverse(sub_shape)

    for shape in slide.shapes:
        traverse(shape)

    return charts

def generate_ppt_from_wrapped_json(wrapped_data, template_path="./power-point-template/main.pptx"):
    data = wrapped_data.get("response", {})
    reviews = data.get("reviews", [])
    sentiment = data.get("analysis", {}).get("sentiment_summary", {})

    avg_rating = round(sum(r["rating"] for r in reviews) / len(reviews), 1) if reviews else 0
    recommend_percent = round(
        len([r for r in reviews if r.get("recommendation")]) / len(reviews) * 100
    ) if reviews else 0

    prs = Presentation(template_path)
    slides = prs.slides

    try:
        def list_text_boxes(presentation, slide_num):
            slide = presentation.slides[slide_num - 1]
            text_boxes = []
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text.strip():
                    text_boxes.append(shape.text.strip())
            return text_boxes

        #  TODO Tes: ambil semua text box dari slide ke-...
        for idx, text in enumerate(list_text_boxes(prs, 7), 1):
            print(f"Text Box {idx}: {text}")
        chart = None
        
        for slide in slides:
            for shape in slide.shapes:
                chart = find_chart(shape)
                if chart:
                    print("✅ Chart found:", chart.chart_type)
                    break
            if chart:
                break
            
            print("❌ Chart not found at all.")
            
        
        # TODO SLIDE ke 2
        slide = prs.slides[1] 
        
        text_shapes = [shape for shape in slide.shapes if shape.has_text_frame and shape.text.strip()]

        # Pastikan jumlah shape cukup
        p_title = text_shapes[2].text_frame.paragraphs[0]
        if p_title.runs:
            p_title.runs[0].text = data["title"]

        p_type = text_shapes[1].text_frame.paragraphs[0]
        if p_type.runs:
            p_type.runs[0].text = data["type"]
            
        
         # TODO SLIDE ke 3
        slide = prs.slides[2] 
        text_shapes = [shape for shape in slide.shapes if shape.has_text_frame and shape.text.strip()]
        p_price = text_shapes[0].text_frame.paragraphs[0]
        p_total_review = text_shapes[3].text_frame.paragraphs[0]
        p_rating = text_shapes[5].text_frame.paragraphs[0]
        p_recommendation = text_shapes[7].text_frame.paragraphs[0]

        if p_price.runs:
            p_price.runs[0].text = data["price"]
        if p_total_review.runs:
            p_total_review.runs[0].text = str(data["total_review"])
        if p_rating.runs:
            p_rating.runs[0].text = f"{avg_rating}/5"
        if p_recommendation.runs:
            p_recommendation.runs[0].text = f"{recommend_percent}%"
            
        # TODO SLIDE ke 4
        slide = prs.slides[3] 
        text_shapes = [shape for shape in slide.shapes if shape.has_text_frame and shape.text.strip()]
        p_pros = text_shapes[0].text_frame.paragraphs[0]
        p_cons = text_shapes[1].text_frame.paragraphs[0]   
        
        top_pros_data = data["analysis"].get("top_pros", [])[:10]
        top_cons_data = data["analysis"].get("top_cons", [])[:10]

        top_pros = [f"{phrase} ({count})" for phrase, count in top_pros_data]
        top_cons = [f"{phrase} ({count})" for phrase, count in top_cons_data]
         
        if p_pros.runs:
            p_pros.runs[0].text = ", ".join(top_pros)
        if p_cons.runs:
            p_cons.runs[0].text = ", ".join(top_cons)  
 
        chart_data = CategoryChartData()
        chart_data.categories = ["Positive", "Neutral", "Negative"]
        chart_data.add_series("Sentiment", [
            sentiment.get("positive", 0),
            sentiment.get("neutral", 0),
            sentiment.get("negative", 0),
        ])

        chart.replace_data(chart_data)
        
        

        if chart.chart_type == XL_CHART_TYPE.PIE:
            chart.replace_data(chart_data)
            print("✅ Chart data successfully replaced.")
        else:
            print(f"⚠️ Chart type {chart.chart_type} cannot be replaced automatically.")
            
        # TODO SLIDE ke 5
        slide = prs.slides[4]
        text_shapes = [shape for shape in slide.shapes if shape.has_text_frame and shape.text.strip()]
    

        p_review_positive = text_shapes[0].text_frame.paragraphs[0]
        p_username_positive = text_shapes[1].text_frame.paragraphs[0]
        p_review_negative = text_shapes[3].text_frame.paragraphs[0]
        p_username_negative = text_shapes[4].text_frame.paragraphs[0]

        # Ambil keyword pros/cons
        top_pro_keywords = [p[0].lower() for p in data["analysis"].get("top_pros", [])]
        top_con_keywords = [c[0].lower() for c in data["analysis"].get("top_cons", [])]

        # Temukan review positif yang mengandung top_pro
        positive_review = next((
            r for r in reviews 
            if r["rating"] >= 4 and any(kw in r["body"].lower() for kw in top_pro_keywords)
        ), None)

        # Jika tidak ketemu, fallback ke review positif mana saja
        if not positive_review:
            positive_review = next((r for r in reviews if r["rating"] >= 4), None)

        # Temukan review negatif yang mengandung top_con
        negative_review = next((
            r for r in reviews 
            if r["rating"] <= 2 and any(kw in r["body"].lower() for kw in top_con_keywords)
        ), None)

        # Jika tidak ketemu, fallback ke review negatif mana saja
        if not negative_review:
            negative_review = next((r for r in reviews if r["rating"] <= 2), None)

        # Isi ke slide jika ditemukan
        if positive_review:
            if p_review_positive.runs:
                p_review_positive.runs[0].text = positive_review["body"]
            if p_username_positive.runs:
                p_username_positive.runs[0].text = f"— {positive_review['author']}"

        if negative_review:
            if p_review_negative.runs:
                p_review_negative.runs[0].text = negative_review["body"]
            if p_username_negative.runs:
                p_username_negative.runs[0].text = f"— {negative_review['author']}"
                
        # TODO SLIDE ke 6
        slide = prs.slides[5]
        text_shapes = [shape for shape in slide.shapes if shape.has_text_frame and shape.text.strip()]
        charts = find_charts_in_slide(slide)
        if "analysis" in data and "rating_insights" in data["analysis"]:
            paragraphs = [
                text_shapes[1].text_frame.paragraphs[0],
                text_shapes[2].text_frame.paragraphs[0],
                text_shapes[3].text_frame.paragraphs[0]
            ]

            insights = data["analysis"]["rating_insights"]

            for para, insight in zip(paragraphs, insights):
                if para.runs:
                    sentiment = insight["sentiment"]
                    emoji = "👍" if sentiment == "positive" else "😑" if sentiment == "neutral" else "👎"
                    para.runs[0].text = f'{emoji}{insight["message"]}'
        else:
            print("Key 'rating_insights' not found in data['analysis']")



        if "analysis" in data and "rating_distribution" in data["analysis"]:
            rating_data = data["analysis"]["rating_distribution"]

            # Bangun chart data baru
            chart_data = CategoryChartData()
            categories = [f"{entry['key']}★" for entry in rating_data]
            values = [entry["data"] for entry in rating_data]
            chart_data.categories = categories
            chart_data.add_series("Ratings", values)
            

            # Replace chart di slide
            for idx, chart in enumerate(charts, 1):
                print(f"Chart slide 6 {idx}: type = {chart.chart_type}")
                try:
                    chart.replace_data(chart_data)
                    print("✅ Chart data berhasil diupdate.")
                except Exception as e:
                    print(f"❌ Gagal mengganti chart data: {e}")
        else:
            print("Key 'rating_distribution' not found in data['analysis']")
            
        # TODO SLIDE ke 7
        slide = prs.slides[6]  # Slide ke-7 index-nya 6
        text_shapes = [shape for shape in slide.shapes if shape.has_text_frame and shape.text.strip()]
        charts = find_charts_in_slide(slide)

        # Masukkan duration_insights ke teks
        if "analysis" in data and "duration_insights" in data["analysis"]:
            paragraphs = [
                text_shapes[1].text_frame.paragraphs[0],
                text_shapes[2].text_frame.paragraphs[0],
                text_shapes[3].text_frame.paragraphs[0]
            ]

            insights = data["analysis"]["duration_insights"]

            for para, insight in zip(paragraphs, insights):
                if para.runs:
                    sentiment = insight["sentiment"]
                    emoji = "👍" if sentiment == "positive" else "😑" if sentiment == "neutral" else "👎"
                    para.runs[0].text = f'{emoji} {insight["message"]}'
        else:
            print("Key 'duration_insights' not found in data['analysis']")

        # Masukkan duration_distribution ke chart
        if "analysis" in data and "duration_distribution" in data["analysis"]:
            duration_data = data["analysis"]["duration_distribution"]

            # Bangun chart data baru
            chart_data = CategoryChartData()
            categories = [entry["key"] for entry in duration_data]
            values = [entry["data"] for entry in duration_data]
            chart_data.categories = categories
            chart_data.add_series("Durations", values)

            # Replace chart
            for idx, chart in enumerate(charts, 1):
                print(f"Chart slide 7 {idx}: type = {chart.chart_type}")
                try:
                    chart.replace_data(chart_data)
                    chart.has_title = True
                    chart.chart_title.text_frame.text = "Duration of Use Distribution"
                    print("✅ Duration chart successfully updated.")
                except Exception as e:
                    print(f"❌ Failed to change duration chart: {e}")
        else:
            print("Key 'duration_distribution' not found in data['analysis']")

        # TODO SlIDE ke 8 - Monthly Trend Line Chart
        slide = prs.slides[7]  # Slide ke-8, index ke-7
        text_shapes = [shape for shape in slide.shapes if shape.has_text_frame and shape.text.strip()]
        charts = find_charts_in_slide(slide)

        # Masukkan trend_insights ke textbox
        if "analysis" in data and "trend_insights" in data["analysis"]:
            paragraphs = [
                text_shapes[1].text_frame.paragraphs[0],
                text_shapes[2].text_frame.paragraphs[0],
                text_shapes[3].text_frame.paragraphs[0]
            ]

            insights = data["analysis"]["trend_insights"]

            for para, insight in zip(paragraphs, insights):
                if para.runs:
                    sentiment = insight["sentiment"]
                    emoji = "✅" if sentiment == "positive" else "😑" if sentiment == "neutral" else "⚠️"
                    para.runs[0].text = f'{emoji} {insight["message"]}'
        else:
            print("Key 'trend_insights' not found in data['analysis']")

        # Masukkan trend_distribution ke dalam line chart
        if "analysis" in data and "trend_distribution" in data["analysis"]:
            trend_data = data["analysis"]["trend_distribution"]

            # Konversi nama bulan menjadi angka (urutan 0–11)
            month_to_index = {
                "January": 0, "February": 1, "March": 2, "April": 3,
                "May": 4, "June": 5, "July": 6, "August": 7,
                "September": 8, "October": 9, "November": 10, "December": 11
            }

            # Susun chart data dalam urutan bulan
            sorted_data = sorted(
                trend_data, key=lambda x: month_to_index.get(x["key"], 99)
            )
            categories = [entry["key"] for entry in sorted_data]
            values = [entry["data"] for entry in sorted_data]

            chart_data = CategoryChartData()
            chart_data.categories = categories
            chart_data.add_series("Monthly Ratings", values)

            # Ganti chart datanya
            for idx, chart in enumerate(charts, 1):
                print(f"Chart slide 8 (monthly trend) {idx}: type = {chart.chart_type}")
                try:
                    chart.replace_data(chart_data)
                    chart.has_title = True
                    chart.chart_title.text_frame.text = "Monthly Review Volume"
                    print("✅ Chart monthly trend berhasil diupdate.")
                except Exception as e:
                    print(f"❌ Gagal mengganti chart monthly trend: {e}")
        else:
            print("Key 'trend_distribution' not found in data['analysis']")

        
        
        pptx_io = io.BytesIO()
        prs.save(pptx_io)
        pptx_io.seek(0)
        return pptx_io, None
    except Exception as e:
        return None, str(e)